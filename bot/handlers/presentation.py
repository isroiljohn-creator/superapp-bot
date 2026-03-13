"""Presentation generator — Gemini AI generates slide content, python-pptx renders PPTX."""
import asyncio
import json
import logging
import os
import tempfile

from aiogram import Router, F
from aiogram.types import Message, FSInputFile
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup

from bot.config import settings
from bot.locales import uz
from bot.keyboards.buttons import ai_workers_reply_keyboard
from db.database import async_session
from services.token_service import (
    get_tokens_async, has_enough_async, spend_tokens_async, add_tokens_async,
    PRES_COST,
)

router = Router(name="presentation")
logger = logging.getLogger("presentation")


class PresentationStates(StatesGroup):
    waiting_topic = State()


PRES_SYSTEM_PROMPT = (
    "Sen professional prezentatsiya yaratuvchisan. Foydalanuvchi mavzu beradi, "
    "sen aniq, chiroyli va ma'lumotli slaydlar tuzasan.\n\n"
    "FAQAT JSON formatda javob ber, boshqa hech narsa qo'shma.\n"
    "Format:\n"
    '{"title": "Prezentatsiya nomi", "slides": [\n'
    '  {"title": "Slayd sarlavhasi", "bullets": ["Nuqta 1", "Nuqta 2", "Nuqta 3"]},\n'
    '  ...\n'
    "]}\n\n"
    "Qoidalar:\n"
    "- 6-10 ta slayd yarat\n"
    "- Har bir slaydda 3-5 ta bullet point bo'lsin\n"
    "- O'zbek tilida (lotin alifbosida) yoz\n"
    "- Birinchi slayd — sarlavha, oxirgi slayd — xulosa bo'lsin\n"
    "- Professional va aniq mazmun yoz"
)


def _call_gemini(prompt: str) -> str:
    """Call Gemini API for presentation content."""
    import urllib.request
    api_key = settings.GEMINI_API_KEY
    if not api_key:
        raise ValueError("GEMINI_API_KEY sozlanmagan")
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={api_key}"

    body = json.dumps({
        "system_instruction": {"parts": [{"text": PRES_SYSTEM_PROMPT}]},
        "contents": [{"role": "user", "parts": [{"text": prompt}]}],
        "generationConfig": {"temperature": 0.7, "maxOutputTokens": 4096},
    }).encode()

    req = urllib.request.Request(url, data=body, headers={"Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=60) as resp:
        data = json.loads(resp.read())

    candidates = data.get("candidates", [])
    if candidates and candidates[0].get("content", {}).get("parts"):
        return candidates[0]["content"]["parts"][0]["text"]
    raise ValueError("Gemini javob bermadi")


def _build_pptx(slide_data: dict) -> str:
    """Build a PPTX file from slide data dict. Returns temp file path."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    bg_color = RGBColor(0x1A, 0x1A, 0x2E)       # Dark navy
    title_color = RGBColor(0xFF, 0xD7, 0x00)      # Gold
    text_color = RGBColor(0xE8, 0xE8, 0xE8)       # Light gray
    accent_color = RGBColor(0x00, 0xD4, 0xAA)      # Teal

    def _set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    # Title slide
    slide_layout = prs.slide_layouts[6]  # Blank
    title_slide = prs.slides.add_slide(slide_layout)
    _set_bg(title_slide)

    # Title text
    from pptx.util import Inches
    txBox = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "Prezentatsiya")
    p.font.size = Pt(44)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "AI yordamida yaratildi — Nuvi Academy Bot"
    p2.font.size = Pt(18)
    p2.font.color.rgb = text_color
    p2.alignment = PP_ALIGN.CENTER

    # Content slides
    slides = slide_data.get("slides", [])
    for i, s in enumerate(slides):
        content_slide = prs.slides.add_slide(slide_layout)
        _set_bg(content_slide)

        # Slide number
        num_box = content_slide.shapes.add_textbox(Inches(12), Inches(6.8), Inches(1), Inches(0.5))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = f"{i + 1}/{len(slides)}"
        np.font.size = Pt(12)
        np.font.color.rgb = text_color
        np.alignment = PP_ALIGN.RIGHT

        # Accent line
        from pptx.shapes.autoshape import Shape
        line = content_slide.shapes.add_shape(
            1, Inches(0.8), Inches(1.5), Inches(2), Emu(36000)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = accent_color
        line.line.fill.background()

        # Title
        title_box = content_slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = s.get("title", "")
        tp.font.size = Pt(32)
        tp.font.color.rgb = title_color
        tp.font.bold = True

        # Bullets
        bullet_box = content_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10.5), Inches(4.5))
        btf = bullet_box.text_frame
        btf.word_wrap = True
        for j, bullet in enumerate(s.get("bullets", [])):
            if j == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()
            bp.text = f"  •  {bullet}"
            bp.font.size = Pt(20)
            bp.font.color.rgb = text_color
            bp.space_after = Pt(12)

    # Save to temp file
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(path)
    return path


@router.message(F.text == uz.AI_WORKERS_KB_PRES)
async def start_presentation(message: Message, state: FSMContext):
    """Start presentation generation flow."""
    if not settings.GEMINI_API_KEY:
        await message.answer(
            "⚠️ AI xizmatlari hozircha sozlanmagan.\n"
            "Admin GEMINI_API_KEY ni sozlashi kerak.",
            reply_markup=ai_workers_reply_keyboard(),
        )
        return

    async with async_session() as session:
        if not await has_enough_async(session, message.from_user.id, PRES_COST):
            tokens = await get_tokens_async(session, message.from_user.id)
            await message.answer(
                uz.AI_WORKERS_NO_TOKENS.format(needed=PRES_COST, have=tokens),
                parse_mode="HTML",
                reply_markup=ai_workers_reply_keyboard(),
            )
            return

    await message.answer(uz.PRES_INTRO, parse_mode="HTML")
    await state.set_state(PresentationStates.waiting_topic)


@router.message(PresentationStates.waiting_topic, F.text)
async def handle_presentation_topic(message: Message, state: FSMContext):
    """Generate presentation from topic."""
    topic = message.text.strip()

    # Check for menu buttons
    menu_buttons = [
        uz.MENU_BTN_AI_WORKERS, uz.MENU_BTN_FREE_LESSONS, uz.MENU_BTN_CLUB,
        uz.MENU_BTN_COURSE, uz.MENU_BTN_PROFILE, uz.MENU_BTN_BACK,
        uz.AI_WORKERS_KB_IMAGE, uz.AI_WORKERS_KB_COPY, uz.AI_WORKERS_KB_CHAT,
        uz.AI_WORKERS_KB_DAILY, uz.AI_WORKERS_KB_TOPUP, uz.AI_WORKERS_KB_BACK,
        uz.AI_WORKERS_KB_PRES, uz.AI_WORKERS_KB_LYRICS,
    ]
    if topic in menu_buttons:
        await state.clear()
        return

    # Spend tokens
    async with async_session() as session:
        if not await spend_tokens_async(session, message.from_user.id, PRES_COST):
            tokens = await get_tokens_async(session, message.from_user.id)
            await state.clear()
            await message.answer(
                uz.AI_WORKERS_NO_TOKENS.format(needed=PRES_COST, have=tokens),
                parse_mode="HTML",
                reply_markup=ai_workers_reply_keyboard(),
            )
            return
        await session.commit()

    status_msg = await message.answer(uz.PRES_GENERATING, parse_mode="HTML")

    try:
        prompt = f"Mavzu: {topic}\n\n6-10 slaydli professional prezentatsiya yarad."
        raw = await asyncio.to_thread(_call_gemini, prompt)

        # Parse JSON (strip markdown code fences if present)
        clean = raw.strip()
        if clean.startswith("```"):
            clean = clean.split("\n", 1)[1] if "\n" in clean else clean[3:]
        if clean.endswith("```"):
            clean = clean[:-3]
        clean = clean.strip()
        if clean.startswith("json"):
            clean = clean[4:].strip()

        slide_data = json.loads(clean)

        # Build PPTX
        pptx_path = await asyncio.to_thread(_build_pptx, slide_data)

        # Send file
        slide_count = len(slide_data.get("slides", []))
        doc = FSInputFile(pptx_path, filename=f"{topic[:40]}.pptx")

        async with async_session() as session:
            tokens = await get_tokens_async(session, message.from_user.id)

        try:
            await status_msg.delete()
        except Exception:
            pass

        await message.answer_document(
            doc,
            caption=(
                f"📊 <b>Prezentatsiya tayyor!</b>\n\n"
                f"📝 Mavzu: <i>{topic[:60]}</i>\n"
                f"📄 Slaydlar soni: {slide_count}\n\n"
                f"💰 Qolgan balans: {tokens:,} so'm"
            ),
            parse_mode="HTML",
        )

        # Track analytics
        try:
            async with async_session() as session:
                from services.analytics import AnalyticsService
                from services.crm import CRMService
                crm = CRMService(session)
                user = await crm.get_user(message.from_user.id)
                if user:
                    analytics = AnalyticsService(session)
                    await analytics.track(user_id=user.id, event_type="presentation")
                    await session.commit()
        except Exception:
            pass

        # Cleanup temp file
        try:
            os.unlink(pptx_path)
        except Exception:
            pass

    except Exception as e:
        # Refund on error
        async with async_session() as session:
            await add_tokens_async(session, message.from_user.id, PRES_COST)
            await session.commit()
        logger.error(f"Presentation error: {type(e).__name__}: {e}")
        try:
            await status_msg.edit_text(uz.PRES_ERROR)
        except Exception:
            pass

    await state.clear()
